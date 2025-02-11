import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

def pdf_to_ppt(pdf_file, ppt_file):
    # Abre o arquivo PDF
    pdf_document = fitz.open(pdf_file)
    
    # Cria uma nova apresentação PowerPoint
    presentation = Presentation()

    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        
        # Converte a página PDF para uma imagem (pixmap)
        pix = page.get_pixmap()
        
        # Converte o pixmap em um objeto de imagem que pode ser usado no PowerPoint
        img_bytes = pix.tobytes("png")  # Converte para formato PNG
        img_stream = io.BytesIO(img_bytes)  # Cria um stream de bytes para a imagem
        img = Image.open(img_stream)

        # Cria um novo slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Layout em branco
        
        # Obtém as dimensões do slide
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        img_width, img_height = img.size

        # Aumenta a imagem em 50% (multiplicando por 1.5)
        img_width = int(img_width * 2)
        img_height = int(img_height * 2)

        # Redimensiona a imagem para caber no slide
        aspect_ratio = img_width / img_height
        if aspect_ratio > 1:
            img_width = slide_width
            img_height = int(slide_width / aspect_ratio)
        else:
            img_height = slide_height
            img_width = int(slide_height * aspect_ratio)

        # Calcula as posições para centralizar a imagem
        left = (slide_width - img_width) / 2
        top = (slide_height - img_height) / 2

        # Adiciona a imagem ao slide (centralizada)
        slide.shapes.add_picture(img_stream, left, top, width=img_width, height=img_height)

    # Salva a apresentação PowerPoint
    presentation.save(ppt_file)
    print(f"Arquivo PPT salvo em: {ppt_file}")


